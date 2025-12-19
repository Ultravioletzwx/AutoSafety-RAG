"""
实用函数：上传保存、文本清洗、PDF/PPTX 转 Markdown。

依赖：
    pip install streamlit pymupdf python-pptx
    # 新增依赖
    pip install transformers torch pdf2image pillow markdown-it-py beautifulsoup4
"""
from pathlib import Path
from typing import List
import tempfile

import fitz  # PyMuPDF
from pptx import Presentation
from streamlit.runtime.uploaded_file_manager import UploadedFile
from llama_index.core import Document

import config
# 导入新模块
from engines.ocr_by_vlm.local_parser  import parse_pdf_to_markdown
from md_processor import process_markdown


def save_uploaded_file(uploaded_file: UploadedFile, upload_dir: Path) -> Path:
    """将 Streamlit 上传文件落地到本地目录。"""
    upload_dir.mkdir(parents=True, exist_ok=True)
    target_path = upload_dir / uploaded_file.name
    with target_path.open("wb") as f:
        f.write(uploaded_file.getbuffer())
    return target_path


def clean_text(text: str) -> str:
    """基础清洗，移除多余空行。"""
    return "\n".join(line.strip() for line in text.splitlines() if line.strip())


def pdf_to_documents(file_path: Path) -> List[Document]:
    """将 PDF 转为结构化的 Document 列表，使用 MinerU 2.5 模型进行解析。"""
    # 创建临时输出目录
    with tempfile.TemporaryDirectory() as tmp_dir:
        try:
            # 1. 使用 MinerU 2.5 将 PDF 转换为 Markdown
            md_path = parse_pdf_to_markdown(str(file_path), tmp_dir)
            
            # 2. 处理 Markdown 生成结构化的 Document 对象
            docs = process_markdown(md_path)
            
            # 3. 添加文件名元数据
            for doc in docs:
                if "file_name" not in doc.metadata:
                    doc.metadata["file_name"] = file_path.name
            
            return docs
        except Exception as e:
            print(f"PDF解析失败: {e}")
            # 回退到传统解析方式
            print("正在使用传统方式解析...")
            docs: List[Document] = []
            with fitz.open(file_path) as pdf:
                for page_idx, page in enumerate(pdf, start=1):
                    text = page.get_text("text")
                    markdown = f"# 第 {page_idx} 页\n\n{clean_text(text)}"
                    docs.append(
                        Document(
                            text=markdown,
                            metadata={"file_name": file_path.name, "page_number": page_idx},
                        )
                    )
            return docs


def pptx_to_documents(file_path: Path) -> List[Document]:
    """将 PPTX 转为按页切分的 Document 列表，附带页码元数据。"""
    prs = Presentation(file_path)
    docs: List[Document] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        text_block = "\n".join(texts)
        markdown = f"# 第 {slide_idx} 页\n\n{clean_text(text_block)}"
        docs.append(
            Document(
                text=markdown,
                metadata={"file_name": file_path.name, "page_number": slide_idx},
            )
        )
    return docs


def file_to_documents(file_path: Path) -> List[Document]:
    """根据扩展名调度解析器。"""
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return pdf_to_documents(file_path)
    if suffix == ".pptx":
        return pptx_to_documents(file_path)
    raise ValueError(f"暂不支持的文件类型: {suffix}")


if __name__ == "__main__":
    # 简单测试：读取 sample.pdf/pptx 并输出文档数量
    config.ensure_dirs()
    sample_pdf = config.UPLOAD_DIR / "sample.pdf"
    if sample_pdf.exists():
        print(f"PDF 文档数: {len(pdf_to_documents(sample_pdf))}")

